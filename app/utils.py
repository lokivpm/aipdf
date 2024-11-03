import random
import smtplib
from email.mime.text import MIMEText
from app.config import SMTP_SERVER, SMTP_PORT, SMTP_USERNAME, SMTP_PASSWORD, EMAIL_FROM
import csv
import PyPDF2
from pptx import Presentation
from typing import List, Union

#email sent function
def send_reset_email(email: str, reset_session_id: str):
    reset_link = f"http://localhost:3000/reset-password?session_id={reset_session_id}"

    msg = MIMEText(f"Click the following link to reset your password: {reset_link}")
    msg['Subject'] = 'Password Reset Request'
    msg['From'] = EMAIL_FROM
    msg['To'] = email

    try:
        with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as server:
            server.starttls() 
            server.login(SMTP_USERNAME, SMTP_PASSWORD)
            server.sendmail(msg['From'], [msg['To']], msg.as_string())
        return True
    except Exception as e:
        print(f"Failed to send email: {e}")
        return False

#otp sent function
def send_otp_email(email: str, otp: str):
    msg = MIMEText(f"Your OTP is: {otp}")
    msg['Subject'] = 'Your OTP Code'
    msg['From'] = EMAIL_FROM
    msg['To'] = email

    try:
        with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as server:
            server.starttls()
            server.login(SMTP_USERNAME, SMTP_PASSWORD)
            server.sendmail(msg['From'], [msg['To']], msg.as_string())
        return True
    except Exception as e:
        print(f"Failed to send email: {e}")
        return False

def partition_pdf(file) -> List[str]:
    """Extract text from a PDF file."""
    # pdf_reader = PyPDF2.PdfReader(file)
    # content = []
    # for page in pdf_reader.pages:
    #     content.append(page.extract_text())
    # return content


    for page_num, img in enumerate(images, start=1):
        text = pytesseract.image_to_string(img)
        full_text += f"\n--- Page {page_num} ---\n{text}"
        print('fulll text',full_text)
        ontent.append(full_text)
        return content

def partition_pptx(file) -> List[str]:
    """Extract text from a PowerPoint file."""
    presentation = Presentation(file)
    content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    return content

def partition_csv(file) -> List[str]:
    """Extract text from a CSV file."""
    content = []
    csv_reader = csv.reader(file.decode('utf-8').splitlines())
    for row in csv_reader:
        content.append(" | ".join(row))  
    return content        



